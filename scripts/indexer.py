# -*- coding: utf-8 -*-
# 엑셀은 셀 좌표 정확 추출 + "타깃 셀 주변 30x20" HTML 미리보기 사전생성
# PDF/PPTX/TXT/HWPX도 함께 인덱싱해서 site/index.json에 저장
import re, json, zipfile, traceback
from pathlib import Path
from html import escape

ROOT = Path(__file__).resolve().parents[1]
OUT_DIR = ROOT / "site"
OUT_DIR.mkdir(parents=True, exist_ok=True)
OUT_FILE = OUT_DIR / "index.json"
PREV_DIR = OUT_DIR / "previews"
PREV_DIR.mkdir(parents=True, exist_ok=True)

SUP = {".pdf":"pdf",".pptx":"pptx",".xlsx":"xlsx",".hwpx":"hwpx",".txt":"txt"}
SCAN_DIRS = [ROOT/"site/files", ROOT/"files"]

LOG = []
def log(s): LOG.append(s)

def rel_link(path:Path):
    """Pages에서 접근 가능한 상대경로 반환"""
    if "site" in path.parts:
        rel = Path(*path.parts[path.parts.index("site")+1:])
        return f"./{rel.as_posix()}"
    rel = path.relative_to(ROOT)
    return f"./{rel.as_posix()}"

# ---------- PDF ----------
def extract_pdf(p:Path):
    try:
        from pypdf import PdfReader
        r = PdfReader(str(p))
        for i,pg in enumerate(r.pages):
            t = (pg.extract_text() or "").strip()
            if t:
                yield {"page":i+1,"text":re.sub(r"\s+"," ",t)[:2000]}
    except Exception as e:
        log(f"ERR_PDF {p} :: {e}")

# ---------- PPTX ----------
def extract_pptx(p:Path):
    try:
        from pptx import Presentation
        prs = Presentation(str(p))
        for i,slide in enumerate(prs.slides):
            buf=[]
            for shp in slide.shapes:
                if hasattr(shp,"text") and shp.text:
                    buf.append(shp.text)
            t="\n".join(buf).strip()
            if t:
                yield {"slide":i+1,"text":re.sub(r"\s+"," ",t)[:1500]}
    except Exception as e:
        log(f"ERR_PPTX {p} :: {e}")

# ---------- XLSX (정확 셀 + 미리보기 HTML 생성) ----------
def extract_xlsx(p:Path):
    from openpyxl import load_workbook
    from openpyxl.utils import get_column_letter
    try:
        wb = load_workbook(str(p), data_only=True, read_only=True)
        for s in wb.sheetnames:
            ws = wb[s]
            hits = 0
            # 너무 큰 파일은 앞쪽 N행만 인덱싱(속도/크기 절충)
            for row in ws.iter_rows(min_row=1, max_row=min(ws.max_row, 2000), values_only=False):
                for cell in row:
                    v = cell.value
                    if v is None: continue
                    txt = str(v).strip()
                    if not txt: continue
                    coord = cell.coordinate  # 정확 주소 (예: C12)
                    yield {
                        "sheet": s,
                        "cell": coord,
                        "text": txt[:400],
                        "preview": build_xlsx_preview(ws, s, coord)  # 사전 렌더 HTML
                    }
                    hits += 1
                    if hits >= 2500: break
                if hits >= 2500: break
    except Exception as e:
        log(f"ERR_XLSX {p} :: {e}")

def build_xlsx_preview(ws, sheet_name, cell_addr):
    """타깃 셀 주변 30x20 HTML을 site/previews/에 저장하고 URL 반환"""
    m = re.match(r"^([A-Z]+)(\d+)$", cell_addr, re.I)
    if not m: return ""
    col_letters, row_str = m.group(1).upper(), m.group(2)

    def col_to_idx(col):
        n=0
        for ch in col: n = n*26 + (ord(ch)-64)
        return n-1

    c = col_to_idx(col_letters)
    r = int(row_str)-1

    r0 = max(0, r-15); r1 = min(ws.max_row-1, r+15)
    c0 = max(0, c-10); c1 = min(ws.max_column-1, c+10)

    def enc(x):
        try: return escape(str(x))
        except: return ""

    from openpyxl.utils import get_column_letter
    head_cols = "".join(f"<th>{get_column_letter(ci+1)}</th>" for ci in range(c0, c1+1))
    rows_html=[]
    for ri in range(r0, r1+1):
        row_cells=[f"<td class='rowhead'>{ri+1}</td>"]
        for ci in range(c0, c1+1):
            v = ws.cell(row=ri+1, column=ci+1).value
            cls = "target" if (ri==r and ci==c) else ""
            row_cells.append(f"<td class='{cls}'>{enc(v) if v is not None else ''}</td>")
        rows_html.append("<tr>"+"".join(row_cells)+"</tr>")

    safe_sheet = re.sub(r'[^0-9A-Za-z가-힣_-]+', '_', sheet_name)
    wbname = ws.parent.properties.title or "wb"
    fname = f"{wbname}_{safe_sheet}_{cell_addr}.html"
    path = PREV_DIR / fname

    html = f"""<!doctype html>
<html><head><meta charset="utf-8" />
<style>
  body{{margin:0;font-family:system-ui,-apple-system,Segoe UI,Roboto,Arial,sans-serif}}
  .grid{{overflow:auto;border:1px solid #e6e8ee;border-radius:12px;background:#fff;max-height:60vh}}
  table{{border-collapse:collapse;width:max(600px,100%);font-size:13px}}
  th,td{{border:1px solid #e6e8ee;padding:4px 6px;white-space:nowrap}}
  th{{background:#f1f5f9;position:sticky;top:0;z-index:1}}
  .rowhead{{position:sticky;left:0;background:#f1f5f9;z-index:1}}
  .target{{background:#fde68a;outline:2px solid #f59e0b}}
</style></head>
<body>
<div class="grid">
  <table>
    <thead><tr><th class="rowhead"></th>{head_cols}</tr></thead>
    <tbody>
      {"".join(rows_html)}
    </tbody>
  </table>
</div>
<script>
  const t = document.querySelector('.target');
  if(t) t.scrollIntoView({{block:'center', inline:'center'}});
</script>
</body></html>"""
    path.write_text(html, encoding="utf-8")
    return f"./previews/{fname}"

# ---------- HWPX ----------
def extract_hwpx(p:Path):
    out=[]
    try:
        with zipfile.ZipFile(p) as z:
            for name in z.namelist():
                if name.endswith(".xml"):
                    try:
                        data=z.read(name).decode("utf-8","ignore")
                        text=re.sub(r"<[^>]+>"," ", data)
                        text=re.sub(r"\s+"," ", text).strip()
                        if text: out.append({"text":text[:2000]})
                    except Exception: pass
    except Exception as e:
        log(f"ERR_HWPX {p} :: {e}")
    return out

# ---------- TXT ----------
def extract_txt(p:Path):
    try:
        t = p.read_text(encoding="utf-8", errors="ignore")
        t = re.sub(r"\s+"," ",t).strip()
        return [{"text":t[:2000]}] if t else []
    except Exception as e:
        log(f"ERR_TXT {p} :: {e}")
        return []

# ---------- main ----------
items=[]
found=0
for base in SCAN_DIRS:
    if not base.exists(): continue
    for p in base.rglob("*"):
        if not p.is_file(): continue
        ext = p.suffix.lower()
        if ext not in SUP: continue
        found += 1
        kind = SUP[ext]
        entry_base = {"file": p.name, "fileType": kind, "link": rel_link(p)}
        try:
            if kind=="pdf":
                for chunk in extract_pdf(p):
                    items.append(entry_base | {"snippet":chunk["text"], "page":chunk["page"]})
            elif kind=="pptx":
                for chunk in extract_pptx(p):
                    items.append(entry_base | {"snippet":chunk["text"], "slide":chunk["slide"]})
            elif kind=="xlsx":
                for chunk in extract_xlsx(p):
                    items.append(entry_base | {
                        "snippet":chunk["text"],
                        "sheet":chunk.get("sheet"),
                        "cell":chunk.get("cell"),
                        "preview":chunk.get("preview")  # 사전 생성 HTML
                    })
            elif kind=="hwpx":
                for chunk in extract_hwpx(p):
                    items.append(entry_base | {"snippet":chunk["text"]})
            elif kind=="txt":
                for chunk in extract_txt(p):
                    items.append(entry_base | {"snippet":chunk["text"]})
        except Exception as e:
            log(f"ERR_PROC {p} :: {e}")

OUT_FILE.write_text(json.dumps(items, ensure_ascii=False), encoding="utf-8")
Path("scripts_index.log").write_text("\n".join(LOG), encoding="utf-8")
print(f"WROTE {OUT_FILE} with {len(items)} entries from {found} files")
